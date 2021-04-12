from pptx import Presentation
import glob
import os
import fitz
import docx2txt
import spacy
from spacy.matcher import Matcher
import re
import json
import pandas as pd


class ResumeParser:
    def __init__(self):
        # Trained Model
        self.model = spacy.load('en_core_web_sm')
        self.matcher = Matcher(self.nlp.vocab)
        self.EMAIL_REG = re.compile(r'[a-z0-9.\-+_]+@[a-z0-9.\-+_]+\.[a-z]+')

    def get_name(self, text):
        nlp_text = self.nlp(text)
        pattern = [{'POS': 'PROPN'}, {'POS': 'PROPN'}]
        self.matcher.add('NAME', None, pattern)
        matches = self.matcher(nlp_text)
        for match_id, start, end in matches:
            span = nlp_text[start:end]
            return span.text
        return ""

    def get_skills(self, text):
        text = self.nlp(text)
        noun_chunks = list(text.noun_chunks)
        tokens = [token.text for token in text if not token.is_stop]
        data = pd.read_csv(os.path.join(os.path.dirname(r'Skills\skills.csv'),
                                        'skills.csv'))
        skills = list(data.columns.values)
        skillset = []
        for token in tokens:
            if token.lower() in skills:
                skillset.append(token)
        for token in noun_chunks:
            token = token.text.lower().strip()
            if token in skills:
                skillset.append(token)
        return [i.capitalize() for i in set([i.lower() for i in skillset])]

    def get_experience(self):
        return None

    def get_salary(self):
        return None

    def get_location(self):
        return None

    def get_roles(self):
        return None

    def get_work_details(self):
        return None

    def get_education(self):
        return


def get_emails(text):
    email = re.findall(r"([^@|\s]+@[^@]+\.[^@|\s]+)", text)
    if email:
        try:
            return email[0].split()[0].strip(';')
        except IndexError:
            return None


def get_contact_number(text):
    phone = re.findall(re.compile(r'[\(]?[1-9][0-9 .\-\(\)]{8,}[0-9]'), text)
    if phone:
        number = ''.join(phone[0])
        if len(number) > 10:
            return '+' + number
        else:
            return number


def get_pptx_txt(prs):
    create_file = open('file', 'w')
    create_file.close()
    os.remove("file")
    with open(r"file", "a") as f:
        for file in glob.glob(prs):
            prse = Presentation(file)
            for slide in prse.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if len(shape.text) > 0:
                            f.write("\n")
                        f.write(shape.text)


def get_document_content(file_path):
    file = get_pptx_txt(file_path)
    with open(file, "r") as reading_text_file:
        contents = reading_text_file.read()

    return contents


def get_txt_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    text = ""
    for page in doc:
        text = text + str(page.getText())
    tx = " ".join(text.split('\t'))
    return tx


def docx_to_txt(file):
    text = docx2txt.process(file)
    if text:
        return text.replace('\t', ' ')
    return None


def read_file(file):
    if file.endswith('docx') or file.endswith('doc'):
        text = docx_to_txt(file)
    elif file.endswith('pdf'):
        text = get_txt_from_pdf(file)
    elif file.endswith('pptx'):
        text = get_document_content(file)
    else:
        text = None

    return text


def create_dict(text):
    names = ResumeParser.get_name(text)
    emails = get_emails(text)
    skills = ResumeParser.get_skills(text)
    number = get_contact_number(text)
    details = {
        "name:": names,
        "contact_details": {
            "mobile": number,
            "email": emails
        },
        "skills": skills
    }
    return details


def create_json(file):
    text = read_file(file)
    dict = create_dict(text)
    with open("sample12.json", "w") as outfile:
        json.dump(dict, outfile, indent=4)
    return True

